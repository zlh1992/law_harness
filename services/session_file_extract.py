#!/usr/bin/env python3
"""Extract text from one already-authorized session file.

Authorization and session isolation live in the DSH plugin. This helper only
parses the exact real path passed by that plugin and emits one bounded JSON
document on stdout.
"""

from __future__ import annotations

import json
import pathlib
import subprocess
import sys
from typing import Callable


MAX_EXTRACTED_CHARS = 8_000_000
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl",
    ".yaml", ".yml", ".xml", ".html", ".htm", ".log", ".sql",
    ".js", ".jsx", ".ts", ".tsx", ".py",
}


def bounded(text: str) -> tuple[str, bool]:
    if len(text) <= MAX_EXTRACTED_CHARS:
        return text, False
    return text[:MAX_EXTRACTED_CHARS], True


def read_text(target: pathlib.Path) -> str:
    data = target.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def read_pdf(target: pathlib.Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(target))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def read_docx(target: pathlib.Path) -> str:
    from docx import Document

    document = Document(str(target))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def read_xlsx(target: pathlib.Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(target), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            parts.append("\t".join("" if value is None else str(value) for value in row))
    return "\n".join(parts)


def read_pptx(target: pathlib.Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(target))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def read_with_textutil(target: pathlib.Path) -> str:
    completed = subprocess.run(
        ["/usr/bin/textutil", "-convert", "txt", "-stdout", str(target)],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=60,
    )
    return completed.stdout.decode("utf-8", errors="replace")


def extract(target: pathlib.Path) -> tuple[str, str]:
    extension = target.suffix.lower()
    readers: dict[str, tuple[str, Callable[[pathlib.Path], str]]] = {
        ".pdf": ("pypdf", read_pdf),
        ".docx": ("python-docx", read_docx),
        ".xlsx": ("openpyxl", read_xlsx),
        ".pptx": ("python-pptx", read_pptx),
        ".doc": ("textutil", read_with_textutil),
        ".rtf": ("textutil", read_with_textutil),
        ".odt": ("textutil", read_with_textutil),
    }
    if extension in TEXT_EXTENSIONS:
        return "text", read_text(target)
    if extension in readers:
        extractor, reader = readers[extension]
        return extractor, reader(target)
    raise ValueError(f"unsupported document type: {extension or 'no extension'}")


def main() -> int:
    if len(sys.argv) != 2:
        raise ValueError("expected exactly one file path")
    target = pathlib.Path(sys.argv[1]).resolve(strict=True)
    if not target.is_file():
        raise ValueError("target is not a regular file")
    extractor, text = extract(target)
    content, truncated = bounded(text.replace("\x00", ""))
    print(json.dumps({"ok": True, "extractor": extractor, "content": content, "truncated": truncated}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as error:  # Keep stderr concise for the model-facing tool.
        print(json.dumps({"ok": False, "error": str(error)}, ensure_ascii=False))
        raise SystemExit(1)
